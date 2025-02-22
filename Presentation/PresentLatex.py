

import matplotlib.pyplot as plt
import yfinance as yf
import requests


from Anlysis.Finnancial_Ratios.scraperyahoo import calc_Ratios_with_growth

def generate_line_chart(x,y):
    # Generate a simple line chart

    for line,label in y:
     plt.plot(x, line, label=label)

    plt.xlabel('X Axis')
    plt.ylabel('Y Axis')
    plt.title('Simple Line Chart')
    plt.legend()

    # Save the chart as an image
    plt.savefig('line_chart.png')
    plt.close()

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from SSH.connect_to_slurm import ssh_connect_and_authenticate
import re

# Function to preprocess text and extract sections
def preprocess_text(text):
    """Remove newlines and extra spaces from the input text."""
    return re.sub(r'\s+', ' ', text).strip()

def extract_section(text, start_marker, end_marker=None):
    """Extract a section of text between start_marker and end_marker."""
    try:
        start_index = text.find(start_marker) + len(start_marker)
        if end_marker:
            end_index = text.find(end_marker, start_index)
            return text[start_index:end_index].strip()
        return text[start_index:].strip()
    except Exception as e:
        print(f"Error extracting section: {e}")
        return ""
def create_presentation(ratios, symbol,stock,local_file_path):
    """
    Generates a PowerPoint presentation with the calculated financial ratios.
    """
    prs = Presentation()


    resp = ssh_connect_and_authenticate(local_file_path=local_file_path,question_content="")
    company_name = extract_section(resp, "**Title Slide**: ", "**Company Overview**")
    company_description = extract_section(resp, "**Company Overview**: ", "**Key Financial Insights**")
    key_insights = extract_section(resp, "**Key Financial Insights**: ", "**Risks and Challenges**")
    risks = extract_section(resp, "**Risks and Challenges**: ", "**Actionable Insights**")
    actions = extract_section(resp, "**Actionable Insights**: ", "**Financial Metrics**")
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = f"{company_name} Financial Presentation"
    subtitle.text = f"Analysis for {symbol}"

    # Company Overview Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Company Overview"
    content = slide.placeholders[1]
    content.text = company_description

    # Key Insights Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Financial Insights"
    content = slide.placeholders[1]
    content.text = key_insights

    # Risks and Challenges Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Risks and Challenges"
    content = slide.placeholders[1]
    content.text = risks

    # Actionable Insights Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Actionable Insights"
    content = slide.placeholders[1]
    content.text = actions
    # Financial Ratios Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Financial Ratios"
    content = slide.placeholders[1]
    text = ""
    for ratio, value in ratios.items():
      try:
        text += f"{ratio}: {value:.2f}\n"
      except:
          text += f"{ratio}: None\n"
    content.text = text

    # Function to split text into chunks of 20 words
    def split_text(text, words_per_line=20):
        words = text.split()
        return '\n'.join(
            ' '.join(words[i:i + words_per_line])
            for i in range(0, len(words), words_per_line)
        )

    formatted_text = split_text(text)

    # Set the formatted text and text size
    content.text = formatted_text
    for paragraph in content.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(16)

    # Line Chart Slide
    x = [1, 2, 3, 4]
    y= [([stock.income_stmt.loc['Total Revenue'].iloc[i] for i in range(4)],'Total Revenue'),
        ([stock.balance_sheet.loc['Cash And Cash Equivalents'].iloc[i] for i in range(4)],'Cash And Cash Equivalents'),
        ([stock.cashflow.loc['Free Cash Flow'].iloc[i] for i in range(4)], 'Free Cash Flow' )]

    generate_line_chart(x,y)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Line Chart"
    chart_path = "line_chart.png"
    slide.shapes.add_picture(chart_path, Inches(1), Inches(1), width=Inches(5))

    # Save Presentation
    prs.save('financial_presentation.pptx')
    return 'Presentation/financial_presentation.pptx'
def process(symbol):
    image_url = "https://mayafiles.tase.co.il/logos/he-IL/000373.jpg"
    image_path = "company_logo.jpg"
    local_file_path = 'C:/Users/yairb/Downloads/elbit_fs3Q.pdf'
    response = requests.get(image_url)
    with open(image_path, 'wb') as file:
        file.write(response.content)
    # Generate the line chart image


    stock = yf.Ticker(symbol)
    ratios=calc_Ratios_with_growth(stock,symbol)
    # Create the presentation
    return create_presentation(ratios,symbol,stock,local_file_path)
if __name__ == '__main__':
    process('ESLT.TA')
