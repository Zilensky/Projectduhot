

import matplotlib.pyplot as plt
import yfinance as yf
import requests


from Anlysis.Finnancial_Ratios.scraperyahoo import calc_Ratios_with_growth

def generate_line_chart():
    # Generate a simple line chart
    x = [1, 2, 3, 4]
    y1 = [10, 15, 13, 17]
    y2 = [5, 7, 8, 9]

    plt.plot(x, y1, label='Line 1')
    plt.plot(x, y2, label='Line 2')

    plt.xlabel('X Axis')
    plt.ylabel('Y Axis')
    plt.title('Simple Line Chart')
    plt.legend()

    # Save the chart as an image
    plt.savefig('line_chart.png')
    plt.close()

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from SSH.connect_to_slurm import ssh_connect_and_authenticate
def create_presentation(ratios, symbol):
    """
    Generates a PowerPoint presentation with the calculated financial ratios.
    """
    prs = Presentation()

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Yair Company"
    subtitle.text = f"Financial Presentation for {symbol}"

    # Logo Slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = f"Company Logo - {symbol}"
    img_path = "company_logo.jpg"
    slide.shapes.add_picture(img_path, Inches(1), Inches(1), width=Inches(2))
    # Breief Slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = f"Brief - {symbol}"
    left = Inches(1)
    top = Inches(3)
    width = Inches(6)
    height = Inches(1)
    question_content = (

        '''Generate a concise briefing that introduces the company to someone who is unfamiliar with it,Based solely on the content of the provided financial statement. The briefing should focus on key aspects such as the company’s operations, products (if mentioned), its history and the company today. don't include financial data. The briefing must be structured to provide clear, factual information without including any additional commentary, explanations, or interpretations. The briefing should be at least 200 words in length.\n 
        ''')

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = ssh_connect_and_authenticate(local_file_path='C:/Users/yairb/Downloads/P1636964-00.pdf',question_content=question_content)


    # Financial Ratios Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Financial Ratios"
    content = slide.placeholders[1]
    text = ""
    for ratio, value in ratios.items():
      try:
        text += f"{ratio}: {value:.2f}\n"
      except:
          text += f"{ratio}: None\n"
    content.text = text

    # Line Chart Slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Line Chart"
    chart_path = "line_chart.png"
    slide.shapes.add_picture(chart_path, Inches(1), Inches(1), width=Inches(5))

    # Save Presentation
    prs.save('financial_presentation.pptx')

image_url = "https://mayafiles.tase.co.il/logos/he-IL/000373.jpg"
image_path = "company_logo.jpg"
response = requests.get(image_url)
with open(image_path, 'wb') as file:
    file.write(response.content)
# Generate the line chart image
generate_line_chart()
symbol='BABA'
stock = yf.Ticker(symbol)
ratios=calc_Ratios_with_growth(stock,symbol)
# Create the presentation
create_presentation(ratios,symbol)

